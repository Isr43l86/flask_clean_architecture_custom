import io
from abc import ABC
from datetime import datetime

import pytz
from pptx import Presentation
from pptx.util import Pt

from ...constants import app_constants


def delete_slides(ppt: Presentation, sections_to_delete: list):
    for index, section in enumerate(sections_to_delete):
        xml_slides = ppt.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[section['index'] - index])


class PresentationGeneratorImpl(ABC):
    @staticmethod
    def generate_ppt_document(data: dict, path: str, sections_to_delete: list):
        new_dict = {f"[{key}]": str(value).upper() for key, value in data.items()}
        new_dict['[currentDate]'] = str(datetime.now(pytz.timezone(app_constants.APP_TIMEZONE)).strftime('%m-%d-%Y'))
        new_dict['[projectVersion]'] = '1'

        ppt = Presentation(path)

        for slide in ppt.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for key, value in new_dict.items():
                        if key in paragraph.text:
                            paragraph.text = paragraph.text.replace(key, value)
                            for run in paragraph.runs:
                                if value in run.text:
                                    run.font.size = Pt(30)

        delete_slides(ppt, sections_to_delete)

        file_stream = io.BytesIO()
        ppt.save(file_stream)
        file_stream.seek(0)

        filename = f"{str(new_dict['[currentDate]']).replace('-', '_')}_{str(new_dict['[projectTitle]']).replace(' ', '_')}_v1.{new_dict['[projectVersion]']}.pptx"
        return file_stream, filename
