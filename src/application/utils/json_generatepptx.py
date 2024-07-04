# from flask import request
from pptx import Presentation
from pptx.util import Inches

# from flask import request

# "slide_theme": 0 o 21 con la plantilla actual,
SLIDE_DATA = {
    "slide_company_name": "kruger corp",
    "slide_theme": 21,
    "slide_title": "Requerimientos técnicos",
    "slide_content": [
        {
            "label": "Frontend",
            "front_tech": ["flutter", "Firebase para manejo de notificaciones", "Consumo servicios", "Desarrollos particulares"]
        },
        {
            "label": "Backend",
            "backend_tech": ["NodeJS", "Servicios .NET", "SQL"]
        },
        {
            "label": "Servers",
            "servers_tech": ["sql", "aplicaciones", "cloud"]
        }
    ]
}


def update_ppt_service():
    root_path = 'D:/kruger/krugerAI/'
    template_path = root_path + 'nueva_plantilla.pptx'
    save_path = root_path + 'pruebas/new_ppt04.pptx'
    prs = Presentation(template_path)

    # Datos recibidos del objeto SLIDE_DATA
    slide_theme = SLIDE_DATA["slide_theme"]
    slide_title = SLIDE_DATA["slide_title"]
    slide_content = SLIDE_DATA["slide_content"]
    slide_company_name = SLIDE_DATA["slide_company_name"]

    # Insertar título de la presentación en la primera diapositiva
    first_slide = prs.slides[0]

    # Buscar un cuadro de texto específico (basado en tu diseño)
    for shape in first_slide.shapes:
        if not shape.has_text_frame:
            continue
        if "initial_slide_title" in shape.text:
            shape.text = slide_company_name
            break
    else:
        raise AttributeError(
            "No se encontró el cuadro de texto para el título en la primera diapositiva.")

    # Insertar una nueva diapositiva con el layout especificado en slide_theme
    slide_layout = prs.slide_layouts[slide_theme]
    new_slide = prs.slides.add_slide(slide_layout)

    # Mover la nueva diapositiva a la posición 5 (índice 4)
    new_slide_position = 4
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slide_element = slides[-1]  # La nueva diapositiva se agrega al final
    xml_slides.remove(new_slide_element)
    xml_slides.insert(new_slide_position, new_slide_element)

    # Agregar título a la nueva diapositiva

    if (slide_theme == 0):
        if new_slide.shapes.title:
            new_slide.shapes.title.text = slide_title
        else:
            left = Inches(0.5)
            top = Inches(1)
            width = Inches(8.5)
            height = Inches(1)
            textbox = new_slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = slide_title
    else:

        if len(new_slide.placeholders) > 1:
            subtitle_placeholder = new_slide.placeholders[1]
            subtitle_placeholder.text = slide_title
        else:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8.5)
            height = Inches(1)
            textbox = new_slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = slide_title

    # Agregar contenido de slide_content a la nueva diapositiva
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(4)
    textbox = new_slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    for content in slide_content:
        label = content.get('label', '')
        # Obtener la lista de tecnologías
        tech_list = content.get(list(content.keys())[1], [])

        # Agregar el label como párrafo
        p = text_frame.add_paragraph()
        p.text = label
        p.bold = True

        # Agregar cada tecnología de la lista
        for tech in tech_list:
            p = text_frame.add_paragraph()
            p.text = f"- {tech}"

    prs.save(save_path)

    return save_path
