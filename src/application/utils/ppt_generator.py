from pptx import Presentation
from pptx.util import Inches
import os


def load_presentation(template_path):
    return Presentation(template_path)


def save_presentation(prs, save_path):
    prs.save(save_path)


def set_first_slide_title(prs, title):
    first_slide = prs.slides[0]
    for shape in first_slide.shapes:
        if not shape.has_text_frame:
            continue
        if "initial_slide_title" in shape.text:
            shape.text = title
            return
    raise AttributeError(
        "No se encontró el cuadro de texto para el título en la primera diapositiva."
    )


def add_slide(prs, slide_theme, slide_title, slide_content, position):
    slide_layout = prs.slide_layouts[slide_theme]
    new_slide = prs.slides.add_slide(slide_layout)

    # Mover la nueva diapositiva a la posición especificada
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slide_element = slides[-1]
    xml_slides.remove(new_slide_element)
    xml_slides.insert(position, new_slide_element)

    set_slide_title(new_slide, slide_theme, slide_title)
    if slide_content:
        set_slide_content(new_slide, slide_content)


def set_slide_title(slide, slide_theme, slide_title):
    if slide_theme == 8:
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
        else:
            add_textbox(slide, Inches(0.5), Inches(
                1), Inches(8.5), Inches(1), slide_title)
    else:
        if len(slide.placeholders) > 1:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = slide_title
        else:
            add_textbox(slide, Inches(1), Inches(1.5),
                        Inches(8.5), Inches(1), slide_title)


def set_slide_content(slide, slide_content):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    for content in slide_content:
        label = content.get('label', '')
        tech_list = content.get(list(content.keys())[1], [])

        p = text_frame.add_paragraph()
        p.text = label
        p.bold = True

        for tech in tech_list:
            p = text_frame.add_paragraph()
            p.text = f"- {tech}"


def add_textbox(slide, left, top, width, height, text):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text


def update_ppt_service(SLIDE_DATA):
    slide_company_name = SLIDE_DATA["slide_company_name"]
    array_slides = SLIDE_DATA["array_slides"]

    root_path = 'D:/kruger/krugerAI/'
    template_path = os.path.join(root_path, 'nueva_plantilla.pptx')
    save_path = os.path.join(root_path, 'pruebas',
                             f'propuesta_{slide_company_name}.pptx')

    prs = load_presentation(template_path)
    set_first_slide_title(prs, slide_company_name)

    for i, slide_data in enumerate(array_slides):
        slide_theme = slide_data["slide_theme"]
        slide_title = slide_data.get("slide_title", "")
        slide_content = slide_data.get("slide_content", [])
        add_slide(prs, slide_theme, slide_title, slide_content, 1 + i)

    save_presentation(prs, save_path)
    return save_path
