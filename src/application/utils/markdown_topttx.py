from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup


texto_markdown = """
6. **Timeline:**
   - **Project Phases:**
     - Planning: 10 days
     - Requirements Analysis: 15 days
     - System Design: 20 days
     - Development: 40 days
     - Testing: 20 days
     - Implementation: 10 days
     - Maintenance: Ongoing
   - **Total Duration:** 115 working days

   - **Project Team:**
     - 1 Project Manager: 960 hours
     - 1 Business Analyst: 160 hours
     - 1 Software Architect: 160 hours
     - 2 Backend Developers: 640 hours (320 hours each)
     - 2 Frontend Developers: 640 hours (320 hours each)
     - 2 Testers: 480 hours (240 hours each)
     - Contingency roles: Technical Support as needed
"""


def markdownToPPTX(markdown_text, text_frame):
    mardownLines = markdown_text.split("\n")

    for line in mardownLines:
        if '**' in line and '-' not in line:
            run = text_frame.add_paragraph().add_run()
            run.text = line.replace('**', '')
            run.font.bold = True
        elif '-' in line and '**' not in line:
            run = text_frame.add_paragraph().add_run()
            run.text = '• ' + line.replace('-', '')
            run.font.bold = False
        elif '**' in line and '*' in line:
            run = text_frame.add_paragraph().add_run()
            run.text = '-' + line.replace('**', '').replace('-', '')
            run.font.bold = True


def update_ppt_service():
    root_path = 'D:/kruger/krugerAI/'
    template_path = root_path + 'nueva_plantilla.pptx'
    save_path = root_path + 'pruebas/new_ppt04.pptx'
    prs = Presentation(template_path)

    # html_text = markdown.markdown(texto_markdown)

    # Valores por defecto
    slide_theme = 0
    slide_title = "nueva slide desde python"
    slide_company_name = "kruger desde python"

    # Insertar título de la presentación en la primera diapositiva
    first_slide = prs.slides[0]

    # Buscar un cuadro de texto específico (basado en tu diseño)
    for shape in first_slide.shapes:
        if not shape.has_text_frame:
            continue
        if "initial_slide_title" in shape.text:
            shape.text = slide_company_name
            break
    else:
        raise AttributeError(
            "No se encontró el cuadro de texto para el título en la primera diapositiva.")

    # Insertar una nueva diapositiva con el layout especificado en slide_theme
    slide_layout = prs.slide_layouts[slide_theme]
    new_slide = prs.slides.add_slide(slide_layout)

    # Mover la nueva diapositiva a la posición 5 (índice 4)
    new_slide_position = 4
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slide_element = slides[-1]  # La nueva diapositiva se agrega al final
    xml_slides.remove(new_slide_element)
    xml_slides.insert(new_slide_position, new_slide_element)

    # Agregar título a la nueva diapositiva
    if new_slide.shapes.title:
        new_slide.shapes.title.text = slide_title
    else:
        left = Inches(0.5)
        top = Inches(1)
        width = Inches(8.5)
        height = Inches(1)
        textbox = new_slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = slide_title

    # Agregar contenido de slide_content a la nueva diapositiva
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(4)
    textbox = new_slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    markdownToPPTX(texto_markdown, text_frame)

    # parse_html_to_pptx(text_frame, html_text)

    prs.save(save_path)

    return save_path
